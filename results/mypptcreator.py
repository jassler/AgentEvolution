from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.action import PP_ACTION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor

import os
from pathlib import Path

import cv2

# img = cv2.imread('combined/coords.png', cv2.IMREAD_UNCHANGED)

# # (1) create a copy of the original:
# overlay = img.copy()
# # (2) draw shapes:
# cv2.circle(overlay, (133, 132), 12, (0, 255, 0), -1)
# cv2.circle(overlay, (166, 132), 12, (0, 255, 0), -1)
# # (3) blend with the original:
# opacity = 0.4
# cv2.addWeighted(overlay, opacity, img, 1 - opacity, 0, img)


beginnings = [
    '0,0,1',
    '0,0.5,0.5',
    '0,1,0',
    '0.1,0.4,0.5',
    '0.1,0.75,0.15',
    '0.15,0.1,0.75',
    '0.15,0.35,0.5',
    '0.2,0.3,0.5',
    '0.333,0.333,0.333',
    '0.4,0.1,0.5',
    '0.5,0,0.5',
    '0.5,0.5,0',
    '0.75,0.15,0.1',
    '1,0,0'
]

beginnings_digits = {
    '0,0,1': [0,0,1],
    '0,0.5,0.5': [0,0.5,0.5],
    '0,1,0': [0,1,0],
    '0.1,0.4,0.5': [0.1,0.4,0.5],
    '0.1,0.75,0.15': [0.1,0.75,0.15],
    '0.15,0.1,0.75': [0.15,0.1,0.75],
    '0.15,0.35,0.5': [0.15,0.35,0.5],
    '0.2,0.3,0.5': [0.2,0.3,0.5],
    '0.333,0.333,0.333': [0.333,0.333,0.333],
    '0.4,0.1,0.5': [0.4,0.1,0.5],
    '0.5,0,0.5': [0.5,0,0.5],
    '0.5,0.5,0': [0.5,0.5,0],
    '0.75,0.15,0.1': [0.75,0.15,0.1],
    '1,0,0': [1,0,0]
}

payoffs = [
    '0-3-,-30-,--0',
    '0+-3,+0-,-3-0',
    '0+-,-30+,-+0',
    '0--,+0-3,--30',
    '0--,+0-,--0',
    '0-+,-0+,--0',
    '0--,+0+,-+0',
    '0+-,+0+,+-0',
    '0++,+0+,--0',
    '0++,-0+,--0',
    '0-3,-03,++0',
    '0+-,-0+,-+0',
    '0+3,-05,+30',
    '03-,30-,++0',
    '0++,+0+,++0',
    '06-4,-305,-30',
    '0++,-03,++0',
    '03-,+0+,3-0',
    '02-,-02,2-0'
]

combined_path = '/Users/felixfritz/Documents/GitHub/AgentEvolution/results/combined'

prs = Presentation()
prs.slide_width = Cm(33.87)
prs.slide_height = Cm(19.05)

blank_slide_layout = prs.slide_layouts[1]

subsections = prs.slide_width // 3

title_slide = prs.slides.add_slide(blank_slide_layout)
title_columns = 5
title_rows = len(payoffs) // title_columns + 1
title_rowheight = prs.slide_height // title_rows
title_columnwidth = prs.slide_width // title_columns
title_shapes = {}
for i, p in enumerate(payoffs):
    column = i % title_columns
    row = i // title_columns
    pic = title_slide.shapes.add_picture(os.path.join(combined_path, 'zeeman', '{}.png'.format(p)),
                            column*title_columnwidth, row*title_rowheight, width=title_columnwidth, height=title_rowheight)
    title_shapes[p] = pic

def create_navigation_bar(chapter_sheet, current_sheet, add_chapter_link = True):
    shapes = current_sheet.shapes
    dimensions = Cm(1)

    backgroundColor = RGBColor(255, 255, 255)
    borderColor = RGBColor(0, 0, 0)

    backgroundColorInactive = RGBColor(122, 122, 122)
    borderColorInactive = RGBColor(122, 122, 122)
    
    rect = shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_HOME, prs.slide_width - dimensions * 2 - Cm(0.5) * 2, prs.slide_height - dimensions - Cm(0.5), dimensions, dimensions)
    rect.line.color.rgb = borderColor
    rect.fill.solid()
    rect.fill.fore_color.rgb = backgroundColor
    rect.click_action.target_slide = title_slide


    rect = shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_RETURN, prs.slide_width - dimensions * 1 - Cm(0.5) * 1, prs.slide_height - dimensions - Cm(0.5), dimensions, dimensions)
    rect.fill.solid()
    
    if add_chapter_link:
        rect.line.color.rgb = borderColor
        rect.fill.fore_color.rgb = backgroundColor
        rect.click_action.target_slide = chapter_sheet
    else:
        rect.line.color.rgb = borderColorInactive
        rect.fill.fore_color.rgb = backgroundColorInactive
    

for p in payoffs:
    beginning_slides = {}
    dot_shapes = {}
    for b in beginnings:
        dot_shapes[b] = []

    img_path = os.path.join(combined_path, p)
    if not os.path.exists(img_path):
        print('{} does not exist'.format(img_path))
        continue
    
    fileamount = len(list(Path(img_path).rglob('*.png')))
    if fileamount != 42:
        print('{} has {} files in it, not 42!'.format(img_path, fileamount))
        continue
    
    # <CHAPTER SLIDE> #
    chapter_slide = prs.slides.add_slide(blank_slide_layout)
    chapter_slide.shapes.title.text = 'Payoff matrix: {}'.format(p)
    chapter_slide.shapes.placeholders[1].text_frame.text = 'LOD average'

    title_shapes[p].click_action.target_slide = chapter_slide

    top = Cm(6)

    # 0--,-0-,--0_0,0.001_avg_LOD
    pic = chapter_slide.shapes.add_picture(os.path.join(combined_path, '{}_0.001,0_avg_LOD.png'.format(p)), 0 * subsections, top, width=subsections)
    pic = chapter_slide.shapes.add_picture(os.path.join(combined_path, '{}_0,0.001_avg_LOD.png'.format(p)), 1 * subsections, top, width=subsections)
    pic = chapter_slide.shapes.add_picture(os.path.join(combined_path, '{}_0.001,0.001_avg_LOD.png'.format(p)), 2 * subsections, top, width=subsections)

    offset = Cm(4.5)
    txBox = chapter_slide.shapes.add_textbox(0 * subsections + offset, top + Cm(8), subsections, Cm(4))
    tf = txBox.text_frame
    tf.text = "g = 0.1%\nm = 0%"

    txBox = chapter_slide.shapes.add_textbox(1 * subsections + offset, top + Cm(8), subsections, Cm(4))
    tf = txBox.text_frame
    tf.text = "g = 0%\nm = 0.1%"

    txBox = chapter_slide.shapes.add_textbox(2 * subsections + offset, top + Cm(8), subsections, Cm(4))
    tf = txBox.text_frame
    tf.text = "g = 0.1%\nm = 0.1%"

    pic = chapter_slide.shapes.add_picture('combined/coords.png', prs.slide_width - Cm(5), Cm(1))
    zeeman_img = chapter_slide.shapes.add_picture('combined/zeeman/{}.png'.format(p), 0, Cm(1), height=pic.height)
    zeeman_img.left = pic.left - zeeman_img.width - Cm(1)
    zeeman_img.click_action.target_slide = title_slide

    create_navigation_bar(chapter_slide, chapter_slide, False)
    # </CHAPTER SLIDE> #

    for b in beginnings:

        slide = prs.slides.add_slide(blank_slide_layout)
        beginning_slides[b] = slide

        shapes = slide.shapes
        shapes.title.text = 'Payoff matrix: {}'.format(p)
        shapes.placeholders[1].text_frame.text = 'Starting point: {}'.format(b)

        top = Cm(6)

        # example: 0,0,1_0,0.001_LOD.png
        pic = shapes.add_picture(os.path.join(img_path, '{}_{}_0.001,0_LOD.png'.format(p, b)), 0 * subsections, top, width=subsections)
        pic = shapes.add_picture(os.path.join(img_path, '{}_{}_0,0.001_LOD.png'.format(p, b)), 1 * subsections, top, width=subsections)
        pic = shapes.add_picture(os.path.join(img_path, '{}_{}_0.001,0.001_LOD.png'.format(p, b)), 2 * subsections, top, width=subsections)

        offset = Cm(4.5)
        txBox = shapes.add_textbox(0 * subsections + offset, top + Cm(8), subsections, Cm(4))
        tf = txBox.text_frame
        tf.text = "g = 0.1%\nm = 0%"

        txBox = shapes.add_textbox(1 * subsections + offset, top + Cm(8), subsections, Cm(4))
        tf = txBox.text_frame
        tf.text = "g = 0%\nm = 0.1%"

        txBox = shapes.add_textbox(2 * subsections + offset, top + Cm(8), subsections, Cm(4))
        tf = txBox.text_frame
        tf.text = "g = 0.1%\nm = 0.1%"

        create_navigation_bar(chapter_slide, slide)

        # image width is 227 x 208 (w x h in pixels)
        # triangle 138 x 120
        # width ratio: 138 / 227
        # height ratio: 120 / 208
        pic = shapes.add_picture('combined/coords.png', prs.slide_width - Cm(5), Cm(1))
        zeeman_img = shapes.add_picture('combined/zeeman/{}.png'.format(p), 0, Cm(1), height=pic.height)
        zeeman_img.left = pic.left - zeeman_img.width - Cm(1)
        zeeman_img.click_action.target_slide = chapter_slide

        triangle_dim = (pic.width * 138 // 227, pic.height * 120 // 230)
        triangle_pos = (pic.left + (pic.width - triangle_dim[0]) // 2 - Cm(0.05), pic.top + (pic.height - triangle_dim[1]) // 2 - Cm(0.1))
        
        oval_dim = Cm(0.25)
        y_offset = -oval_dim // 2 + Cm(0.1)
        x_offset = -oval_dim // 2 + Cm(0.1)

        # shapes.add_shape(MSO_SHAPE.RECTANGLE, pic.left + (pic.width - triangle_dim[0]) // 2, pic.top + (pic.height - triangle_dim[1]) // 2, triangle_dim[0], triangle_dim[1])
        for key, val in beginnings_digits.items():
            x_pos = triangle_pos[0] + (val[2] - val[1]) * (triangle_dim[0] // 2) + triangle_dim[0] // 2
            y_pos = triangle_pos[1] + triangle_dim[1] - val[0] * triangle_dim[1]
            shape = shapes.add_shape(MSO_SHAPE.OVAL, x_pos + x_offset, y_pos + y_offset, oval_dim, oval_dim)
            shape.line.color.rgb = RGBColor(0, 0, 0)

            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)

            if key == b:
                shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
            
            dot_shapes[key].append(shape)
    
    for key, shapes in dot_shapes.items():
        for shape in shapes:
            shape.click_action.target_slide = beginning_slides[key]
        


prs.save('zeemen.pptx')
